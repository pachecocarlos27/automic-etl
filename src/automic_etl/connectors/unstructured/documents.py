"""General document connector for Word, PowerPoint, and other formats."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

import polars as pl

from automic_etl.connectors.base import (
    ConnectorConfig,
    ConnectorType,
    ExtractionResult,
    UnstructuredConnector,
)
from automic_etl.core.exceptions import ExtractionError


@dataclass
class DocumentConfig(ConnectorConfig):
    """Document connector configuration."""

    path: str = ""
    extract_images: bool = False
    extract_tables: bool = True

    def __post_init__(self) -> None:
        self.connector_type = ConnectorType.UNSTRUCTURED


class DocumentConnector(UnstructuredConnector):
    """General document connector supporting Word, PowerPoint, etc."""

    SUPPORTED_EXTENSIONS = {
        ".docx": "word",
        ".doc": "word",
        ".pptx": "powerpoint",
        ".ppt": "powerpoint",
        ".xlsx": "excel",
        ".xls": "excel",
        ".txt": "text",
        ".md": "markdown",
        ".html": "html",
        ".htm": "html",
    }

    def __init__(self, config: DocumentConfig) -> None:
        super().__init__(config)
        self.doc_config = config

    def connect(self) -> None:
        """Verify path exists."""
        self._connected = True
        self.logger.info("Document connector ready", path=self.doc_config.path)

    def disconnect(self) -> None:
        """No persistent connection to close."""
        self._connected = False

    def test_connection(self) -> bool:
        """Test if path is accessible."""
        path = Path(self.doc_config.path)
        return path.exists() or path.parent.exists()

    def extract(
        self,
        query: str | None = None,
        path: str | None = None,
        **kwargs: Any,
    ) -> ExtractionResult:
        """Extract content from document."""
        file_path = Path(path or self.doc_config.path)

        try:
            content = self.extract_content(str(file_path))
            metadata = self.extract_metadata(str(file_path))

            data = {
                "file_path": [str(file_path)],
                "file_type": [file_path.suffix.lower()],
                "text": [content.get("text", "")],
                "word_count": [len(content.get("text", "").split())],
            }

            for key, value in metadata.items():
                if isinstance(value, (str, int, float, bool)):
                    data[f"meta_{key}"] = [value]

            df = pl.DataFrame(data)

            return ExtractionResult(
                data=df,
                row_count=1,
                metadata={
                    "file": str(file_path),
                    "file_type": file_path.suffix.lower(),
                    **metadata,
                },
            )
        except Exception as e:
            raise ExtractionError(
                f"Failed to extract document: {str(e)}",
                source=str(file_path),
            )

    def extract_content(self, source: str | bytes) -> dict[str, Any]:
        """Extract content from document source."""
        if isinstance(source, bytes):
            return {"text": "", "error": "Bytes input not supported"}

        file_path = Path(source)
        extension = file_path.suffix.lower()
        doc_type = self.SUPPORTED_EXTENSIONS.get(extension, "unknown")

        if doc_type == "word":
            return self._extract_word(file_path)
        elif doc_type == "powerpoint":
            return self._extract_powerpoint(file_path)
        elif doc_type == "excel":
            return self._extract_excel(file_path)
        elif doc_type == "text":
            return self._extract_text(file_path)
        elif doc_type == "markdown":
            return self._extract_text(file_path)
        elif doc_type == "html":
            return self._extract_html(file_path)
        else:
            return self._extract_generic(file_path)

    def _extract_word(self, path: Path) -> dict[str, Any]:
        """Extract content from Word document."""
        try:
            from docx import Document
            doc = Document(str(path))

            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(paragraphs)

            tables = []
            if self.doc_config.extract_tables:
                for table in doc.tables:
                    table_data = []
                    for row in table.rows:
                        table_data.append([cell.text for cell in row.cells])
                    tables.append(table_data)

            return {
                "text": text,
                "tables": tables,
                "paragraph_count": len(paragraphs),
            }
        except ImportError:
            return {"text": "", "error": "python-docx not installed"}
        except Exception as e:
            return {"text": "", "error": str(e)}

    def _extract_powerpoint(self, path: Path) -> dict[str, Any]:
        """Extract content from PowerPoint presentation."""
        try:
            from pptx import Presentation
            prs = Presentation(str(path))

            slides_text = []
            for slide in prs.slides:
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content.append(shape.text)
                slides_text.append("\n".join(slide_content))

            return {
                "text": "\n\n---\n\n".join(slides_text),
                "slide_count": len(prs.slides),
            }
        except ImportError:
            return {"text": "", "error": "python-pptx not installed"}
        except Exception as e:
            return {"text": "", "error": str(e)}

    def _extract_excel(self, path: Path) -> dict[str, Any]:
        """Extract content from Excel spreadsheet."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), data_only=True)

            sheets_content = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    rows.append("\t".join(row_text))
                sheets_content.append(f"Sheet: {sheet_name}\n" + "\n".join(rows))

            return {
                "text": "\n\n".join(sheets_content),
                "sheet_count": len(wb.sheetnames),
                "sheets": wb.sheetnames,
            }
        except ImportError:
            return {"text": "", "error": "openpyxl not installed"}
        except Exception as e:
            return {"text": "", "error": str(e)}

    def _extract_text(self, path: Path) -> dict[str, Any]:
        """Extract content from text file."""
        try:
            text = path.read_text(encoding="utf-8")
            return {"text": text, "line_count": len(text.splitlines())}
        except UnicodeDecodeError:
            text = path.read_text(encoding="latin-1")
            return {"text": text, "line_count": len(text.splitlines())}

    def _extract_html(self, path: Path) -> dict[str, Any]:
        """Extract content from HTML file."""
        try:
            from bs4 import BeautifulSoup
            html = path.read_text(encoding="utf-8")
            soup = BeautifulSoup(html, "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            text = soup.get_text(separator="\n", strip=True)
            return {"text": text, "title": soup.title.string if soup.title else ""}
        except ImportError:
            # Fallback without BeautifulSoup
            import re
            html = path.read_text(encoding="utf-8")
            text = re.sub(r"<[^>]+>", " ", html)
            text = re.sub(r"\s+", " ", text).strip()
            return {"text": text}
        except Exception as e:
            return {"text": "", "error": str(e)}

    def _extract_generic(self, path: Path) -> dict[str, Any]:
        """Generic extraction using unstructured library."""
        try:
            from unstructured.partition.auto import partition
            elements = partition(filename=str(path))
            text = "\n\n".join(str(el) for el in elements)
            return {"text": text, "element_count": len(elements)}
        except ImportError:
            return {"text": "", "error": "unstructured not installed"}
        except Exception as e:
            return {"text": "", "error": str(e)}

    def extract_metadata(self, source: str | bytes) -> dict[str, Any]:
        """Extract metadata from document source."""
        if isinstance(source, bytes):
            return {}

        file_path = Path(source)
        stat = file_path.stat()

        metadata = {
            "file_name": file_path.name,
            "file_size": stat.st_size,
            "modified_time": stat.st_mtime,
            "extension": file_path.suffix.lower(),
        }

        # Try to get document-specific metadata
        extension = file_path.suffix.lower()

        if extension in [".docx"]:
            try:
                from docx import Document
                doc = Document(str(file_path))
                core = doc.core_properties
                metadata.update({
                    "author": core.author or "",
                    "title": core.title or "",
                    "created": str(core.created) if core.created else "",
                    "modified": str(core.modified) if core.modified else "",
                })
            except Exception:
                pass

        return metadata
